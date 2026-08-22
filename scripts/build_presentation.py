#!/usr/bin/env python3
"""ISO/TC 154 DIS proposal deck — v2: narrative + SVG diagrams."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

HERE = Path(__file__).resolve().parent.parent
DIA = HERE / "presentations" / "diagrams"

INK = RGBColor(0x0A, 0x16, 0x28)
PAPER = RGBColor(0xFF, 0xFA, 0xF2)
COPPER = RGBColor(0xB8, 0x5A, 0x2A)
CYAN = RGBColor(0x2F, 0x6F, 0x82)
MIST = RGBColor(0x8A, 0x9A, 0xAB)
SLATE = RGBColor(0x24, 0x36, 0x47)
ICE = RGBColor(0xC5, 0xDD, 0xE4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def band(s, color=COPPER, h=0.12):
    from pptx.enum.shapes import MSO_SHAPE
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()


def text(s, x, y, w, h, body, size=18, color=INK, bold=False, mono=False, align=None, leading=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in body.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold
        if mono: p.font.name = "Courier New"
        if align == "center": p.alignment = PP_ALIGN.CENTER
        if leading: p.space_after = Pt(leading)
    return tb


def eyebrow(s, t, y=0.5):
    text(s, 0.9, y, 11.5, 0.4, t, size=14, color=CYAN, mono=True)


def headline(s, t, y=0.92, size=34, color=INK, w=11.5):
    text(s, 0.9, y, w, 1.05, t, size=size, color=color, bold=True)


def kicker(s, t, y=6.35, color=COPPER, size=16):
    text(s, 0.9, y, 11.5, 0.75, t, size=size, color=color, bold=True)


def diagram(s, png, top=1.7, height=5.0):
    from PIL import Image
    p = DIA / png
    iw, ih = Image.open(p).size
    scale = min(11.6 / (iw / 200), height / (ih / 200))
    w, h = (iw / 200) * scale, (ih / 200) * scale
    x = (W - w) / 2
    s.shapes.add_picture(str(p), Inches(x), Inches(top), Inches(w), Inches(h))


def pageno(s, n):
    text(s, 0.9, 7.05, 11.5, 0.35,
         f"CalConnect → ISO/TC 154 · CC/ISO 36010 · {n}",
         size=10, color=MIST)


# ---- 1 · Title ----
s = slide(INK)
band(s, COPPER, 0.2)
text(s, 0.9, 1.35, 11.5, 0.5, "PROPOSAL FOR A DRAFT INTERNATIONAL STANDARD", size=16, color=CYAN, mono=True)
text(s, 0.9, 1.95, 11.5, 2.1, "Lightweight document —\nDocument metamodel", size=52, color=PAPER, bold=True)
text(s, 0.9, 4.25, 11.5, 0.7, "The missing interoperability layer for structured text", size=24, color=ICE)
text(s, 0.9, 5.5, 11.5, 1.2,
     "Proposed by CalConnect (TC VCARD)\nBase text CC/ISO 36010 · proof: 6+ years in production · every claim machine-checked",
     size=15, color=MIST, leading=8)

# ---- 2 · The everyday problem (D1) ----
s = slide(); band(s); eyebrow(s, "THE PROBLEM EVERY TEAM KNOWS")
headline(s, "Same meaning, five languages — and none of them travel")
diagram(s, "d1-meaning-lost.png")
pageno(s, 2)

# ---- 3 · Doesn't this exist? (D2) ----
s = slide(); band(s); eyebrow(s, "“DOESN’T SOMETHING LIKE THIS EXIST?”")
headline(s, "No. Every layer exists — except the one that carries meaning")
diagram(s, "d2-missing-layer.png", height=5.15)
pageno(s, 3)

# ---- 4 · HTML/RTF (D3) ----
s = slide(); band(s); eyebrow(s, "“AREN’T HTML AND RTF ALREADY THAT LAYER?”")
headline(s, "HTML and RTF say how text looks — not what it means")
diagram(s, "d3-not-a-serialization.png", height=4.85)
kicker(s, "The metamodel is what all of them mean. Rendering stays exactly where it is.", y=6.5)
pageno(s, 4)

# ---- 5 · The model (D4) ----
s = slide(); band(s); eyebrow(s, "THE ANSWER")
headline(s, "One small model, deliberately lightweight")
diagram(s, "d4-tiers.png", height=4.85)
pageno(s, 5)

# ---- 6 · Why TC 154 (D5) ----
s = slide(); band(s); eyebrow(s, "WHY TC 154 — AND NOT ANOTHER TC?")
headline(s, "Because TC 154 is where documents are data")
diagram(s, "d5-tc154.png", top=1.65, height=4.9)
pageno(s, 6)

# ---- 7 · Why not exist until now ----
s = slide(); band(s); eyebrow(s, "WHY HAS THIS NOT EXISTED UNTIL NOW?")
headline(s, "Syntaxes competed; consortia standardized the pieces")
why = [
    ("Standards bodies standardized elements, not documents",
     "ISO gave us time (8601), language (639), scripts (15924), references (690). The document carrying them was always someone’s tool, never a standard."),
    ("Format wars had no neutral referee",
     "Markdown vs AsciiDoc vs RST is a tooling debate. A metamodel below all of them was in no vendor’s interest to build."),
    ("It took a standards producer to need it",
     "CalConnect/Metanorma authors standards IN lightweight markup and renders to ISO-grade output — the exact round-trip everyone else gave up on. The model was extracted from production, not invented."),
    ("Now that it is proven, it can be neutral",
     "Six years, 30+ document flavours, executable conformance — ready to hand to ISO as a public good."),
]
y = 1.95
for k, v in why:
    text(s, 0.9, y, 4.6, 0.9, k, size=15.5, bold=True, color=COPPER)
    text(s, 5.8, y, 6.6, 1.15, v, size=13, color=SLATE)
    y += 1.22
pageno(s, 7)

# ---- 8 · Harmonized mappings ----
s = slide(); band(s); eyebrow(s, "HARMONIZED — ANNEX A, NORMATIVE")
headline(s, "Every markup construct has a home: proven, clause by clause")
maps = [
    ("AsciiDoc", "header → bibdata + register\nadmonitions → AdmonitionBlock + type\ncontinuations → relaxed ListItem\npassthroughs → FormattedString\n{attr} → ReferenceToVariable"),
    ("Markdown · GFM · Pandoc", "GFM tables → TableBlock\n{#id .class} → attribute register\n::: divs → register + any container\n$math$ → StemElement (LaTeX)\nraw HTML → FormattedString + format"),
    ("reStructuredText", "directives → register + relaxed content\n|substitution| → ReferenceToVariable\nroles → text specializations\nfield lists → document register\ndanger / error / hint → type extensions"),
]
x = 0.9
for head, body in maps:
    text(s, x, 2.0, 3.8, 0.5, head, size=17, bold=True, color=CYAN)
    text(s, x, 2.55, 3.8, 2.9, body, size=12.5, color=SLATE, mono=True)
    x += 4.0
from pptx.enum.shapes import MSO_SHAPE
box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.15))
box.fill.solid(); box.fill.fore_color.rgb = INK; box.line.fill.background()
text(s, 1.2, 5.64, 10.9, 0.95,
     "Conformance formula — a foreign construct is covered iff it maps to: a construct · a type-specialization · register + relaxed or opaque content.\nAnything unlisted is answerable in one line. That is what “complete superset” means, testably.",
     size=14, color=PAPER, bold=True, leading=4)
pageno(s, 8)

# ---- 9 · Specialize/extend ----
s = slide(); band(s); eyebrow(s, "SPECIALIZE AND EXTEND — ANNEXES B AND C")
headline(s, "Dialects arrive without ever forking the model")
text(s, 0.9, 1.95, 5.5, 0.4, "SPECIALIZE", size=16, bold=True, color=CYAN, mono=True)
text(s, 0.9, 2.42, 5.5, 3.4,
     "· Type values, not parallel classes —\n  admonition kinds, list numeration,\n  stem languages, document classes\n· Attribute overrides: add, remove,\n  retighten (e.g. no hanging paragraphs)\n· Subclassing under the construct roots\n· Constraint specialization:\n  PureTextElement recursion guarantees",
     size=13.5, color=SLATE)
text(s, 7.0, 1.95, 5.4, 0.4, "EXTEND", size=16, bold=True, color=COPPER, mono=True)
text(s, 7.0, 2.42, 5.4, 3.5,
     "· Attribute register: open key-value\n  entries on any construct — dialect\n  keys (checkbox, option), no model change\n· Raw stays lossless: FormattedString,\n  format-qualified blocks\n· Unknown directives (all of Sphinx)\n  decompose into register + relaxed content\n· A document is a block: whole documents\n  compose as child content",
     size=13.5, color=SLATE)
kicker(s, "Open-closed by construction: the basis is defined once and never modified — only specialized upon.", y=6.25)
pageno(s, 9)

# ---- 10 · Conformance you can run (D6) ----
s = slide(); band(s); eyebrow(s, "EXECUTABLE CONFORMANCE — ANNEX D")
headline(s, "Not “trust us” — a standard with a passing test suite")
diagram(s, "d6-conformance.png", height=4.55)
pageno(s, 10)

# ---- 11 · Production proof ----
s = slide(); band(s); eyebrow(s, "EVIDENCE")
headline(s, "Extracted from production, not invented for a committee")
ev = [
    ("30+ document flavours", "Metanorma authors and publishes real standards on this model today (relaton-models, standoc-models consumers, all pinned in lockstep)"),
    ("ISO-grade outputs", "The model renders to ISO/IEC-standard layouts — HTML, PDF, Word — proving the model→presentation direction daily"),
    ("Registry-native", "Localized strings are ISO 24229 spelling systems; romanization schemes carry ISO 24229 system codes as register schemes"),
    ("Open artifacts", "Model sources, grammars, twin instances, gates, atlas site — public at metanorma/basicdoc-models; every merge runs the suite"),
    ("Live model atlas", "metanorma.github.io/basicdoc-models — every plate with full definitions, hyperlinked types, deep links"),
]
y = 1.95
for k, v in ev:
    text(s, 0.9, y, 3.4, 0.6, k, size=15.5, bold=True, color=CYAN)
    text(s, 4.6, y, 7.9, 1.0, v, size=12.5, color=SLATE)
    y += 0.98
pageno(s, 11)

# ---- 12 · The ask ----
s = slide(INK)
band(s, COPPER, 0.2)
text(s, 0.9, 1.15, 11.5, 0.5, "THE ASK", size=16, color=CYAN, mono=True)
text(s, 0.9, 1.7, 11.5, 1.5, "Adopt CC/ISO 36010 as the base text\nfor a Draft International Standard",
     size=38, color=PAPER, bold=True)
asks = [
    "Recognize the lightweight document metamodel as the interchange layer for structured text",
    "Progress it at TC 154: NP → WD → CD → DIS — the base text is complete and tested today",
    "Designate the public repository and its executable conformance suite as the maintenance vehicle",
    "CalConnect (TC VCARD) continues as proposer and editor, liaising with TC 46 — we consume ISO 690 and ISO 5127, we do not duplicate them",
]
y = 3.6
for a in asks:
    text(s, 1.1, y, 11.0, 0.62, "→  " + a, size=16, color=ICE)
    y += 0.74
text(s, 0.9, 6.75, 11.5, 0.4,
     "metanorma.github.io/basicdoc-models · github.com/metanorma/basicdoc-models · CC/ISO 36010",
     size=12, color=MIST, mono=True)

out = HERE / "presentations" / "lightweight-doc-tc154-dis-proposal.pptx"
prs.save(str(out))
print(f"saved: {out} ({len(prs.slides._sldIdLst)} slides)")
